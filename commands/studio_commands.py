"""Async commands for studio resource generation (Project C).

Four independent generators, each a named "agent" run in the worker process
via surreal-commands (same mechanism as podcasts). Each command updates a
single ``StudioArtifact`` record and links it to its command job:

  * generate_report   — ReportWriter:     LLM -> markdown (artifact.content)
  * generate_quiz     — QuizMaster:       LLM -> markdown quiz (artifact.content)
  * generate_mindmap  — MindMapArchitect: LLM -> mermaid  (artifact.content)
  * generate_ppt      — SlideComposer:    LLM -> N slide prompts -> Doubao
                        images -> stitched into one .pptx (one image per slide)
  * generate_video    — VideoDirector:    LLM -> prompt -> Doubao Seedance mp4

Text generation reuses ``provision_langchain_model`` (the same path as
transformations); image/video reuse the Doubao clients from Project A. The PPT
deck is assembled from the generated slide images with python-pptx. Binary
outputs land under ``data/studio/<artifact_id>/``.
"""

import json
import re
import time
from pathlib import Path
from typing import Any, Dict, List, Optional

import httpx
from langchain_core.messages import HumanMessage, SystemMessage
from loguru import logger
from surreal_commands import CommandInput, CommandOutput, command

from open_notebook.ai.provision import provision_langchain_model
from open_notebook.config import DATA_FOLDER
from open_notebook.domain.artifact import StudioArtifact
from open_notebook.utils import clean_thinking_content
from open_notebook.utils.text_utils import extract_text_content

_DOWNLOAD_TIMEOUT = 180.0


def _download(url: str, dest: Path) -> None:
    """Stream a remote file to ``dest``."""
    dest.parent.mkdir(parents=True, exist_ok=True)
    with httpx.stream(
        "GET", url, timeout=_DOWNLOAD_TIMEOUT, follow_redirects=True
    ) as resp:
        resp.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in resp.iter_bytes():
                f.write(chunk)


def _artifact_dir(artifact_id: str) -> Path:
    """Filesystem-safe output dir for an artifact's binary products."""
    safe = str(artifact_id).replace(":", "_")
    return Path(f"{DATA_FOLDER}/studio/{safe}")


async def _run_llm(system_prompt: str, content: str, instructions: Optional[str]) -> str:
    """Run a single LLM completion and return cleaned text output."""
    sys_text = system_prompt
    if instructions:
        sys_text = f"{system_prompt}\n\n额外要求：{instructions}"
    payload = [
        SystemMessage(content=sys_text),
        HumanMessage(content=content or ""),
    ]
    chain = await provision_langchain_model(
        str(payload), None, "transformation", max_tokens=8192
    )
    response = await chain.ainvoke(payload)
    return clean_thinking_content(extract_text_content(response.content))


def _strip_code_fence(text: str) -> str:
    """Strip a leading/trailing ``` fence if the model added one."""
    t = text.strip()
    if t.startswith("```"):
        lines = t.splitlines()
        # drop first fence line (``` or ```lang)
        lines = lines[1:]
        # drop trailing fence line
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        t = "\n".join(lines).strip()
    return t


def repair_mermaid(code: str) -> str:
    """Repair Mermaid syntax errors in mindmap and graph TD to prevent rendering crashes."""
    if not code:
        return code
    
    lines = code.splitlines()
    if not lines:
        return code
    
    # Detect mermaid diagram type
    first_line_stripped = lines[0].strip()
    is_mindmap = first_line_stripped.startswith("mindmap")
    is_graph = first_line_stripped.startswith("graph ") or first_line_stripped.startswith("flowchart ")
    
    if not (is_mindmap or is_graph):
        return code

    repaired_lines = []
    
    if is_mindmap:
        for line in lines:
            stripped = line.strip()
            # Preserve leading spaces (indentation)
            leading_spaces = line[:len(line) - len(line.lstrip())]
            
            # Skip empty lines, 'mindmap' header, config lines, and classDef/style
            if not stripped or stripped.startswith("mindmap") or stripped.startswith("::") or stripped.startswith("classDef") or stripped.startswith("style"):
                repaired_lines.append(line)
                continue
            
            # Check if node text is already wrapped/quoted
            # 1. Enclosed with parens/brackets/quotes: (text), [text], {text}, ("text"), ["text"]
            # 2. explicit node shape with ID: node_id(text), node_id((text)), node_id["text"]
            has_brackets = False
            if (stripped.startswith("(") and stripped.endswith(")")) or \
               (stripped.startswith("[") and stripped.endswith("]")) or \
               (stripped.startswith("{") and stripped.endswith("}")) or \
               (stripped.startswith("\"") and stripped.endswith("\"")):
                has_brackets = True
            elif re.match(r'^[a-zA-Z0-9_\-]+\s*(?:\(\(|\(\(\(|[\(\[\{])', stripped):
                if stripped.endswith(")") or stripped.endswith("]") or stripped.endswith("}"):
                    has_brackets = True
            
            if has_brackets:
                repaired_lines.append(line)
            else:
                # Wrap bare text in (" ") to make it safe for spaces/special characters
                safe_text = stripped.replace('"', '\\"')
                repaired_lines.append(f'{leading_spaces}("{safe_text}")')
                
    elif is_graph:
        node_id_map = {}
        node_counter = 0
        
        def get_node_id(text_content: str) -> str:
            nonlocal node_counter
            text_content = text_content.strip()
            if text_content in node_id_map:
                return node_id_map[text_content]
            
            # If the text is already a clean identifier, reuse it as ID
            if re.match(r'^[a-zA-Z0-9_\-]+$', text_content):
                return text_content
                
            new_id = f"node_{node_counter}"
            node_counter += 1
            node_id_map[text_content] = new_id
            return new_id

        for line in lines:
            stripped = line.strip()
            if not stripped or stripped.startswith("graph ") or stripped.startswith("flowchart "):
                repaired_lines.append(line)
                continue
                
            # Match edge connections (e.g. -->, ---, -.->, ==>)
            connector_pattern = r'(\s*(?:-->|---|-.->|==>|-.->|--\s*[^-]+\s*-->)\s*)'
            parts = re.split(connector_pattern, stripped)
            
            if len(parts) > 1:
                new_parts = []
                for i, part in enumerate(parts):
                    if i % 2 == 0:
                        part_stripped = part.strip()
                        if not part_stripped:
                            new_parts.append("")
                            continue
                        
                        # Already parsed format like ID["text"], ID("text")
                        match = re.match(r'^([a-zA-Z0-9_\-]+)\s*(?:\(\(|\(\(\(|[\(\[\{\"\'])(.*)(?:\)\)|\)\)\)|[\)\]\}\"\'])$', part_stripped)
                        if match:
                            new_parts.append(part_stripped)
                        elif re.match(r'^[a-zA-Z0-9_\-]+$', part_stripped):
                            new_parts.append(part_stripped)
                        else:
                            nid = get_node_id(part_stripped)
                            safe_text = part_stripped.replace('"', '\\"')
                            new_parts.append(f'{nid}["{safe_text}"]')
                    else:
                        new_parts.append(part)
                
                leading_spaces = line[:len(line) - len(line.lstrip())]
                repaired_lines.append(leading_spaces + "".join(new_parts))
            else:
                if not (stripped.startswith("classDef") or stripped.startswith("style") or stripped.startswith("click")):
                    match = re.match(r'^([a-zA-Z0-9_\-]+)\s*(?:\(\(|\(\(\(|[\(\[\{\"\'])(.*)(?:\)\)|\)\)\)|[\)\]\}\"\'])$', stripped)
                    if not match and not re.match(r'^[a-zA-Z0-9_\-]+$', stripped):
                        nid = get_node_id(stripped)
                        safe_text = stripped.replace('"', '\\"')
                        leading_spaces = line[:len(line) - len(line.lstrip())]
                        repaired_lines.append(f'{leading_spaces}{nid}["{safe_text}"]')
                        continue
                repaired_lines.append(line)
                
    return "\n".join(repaired_lines)


# --------------------------------------------------------------------------
# Shared I/O models
# --------------------------------------------------------------------------
class StudioGenerationInput(CommandInput):
    artifact_id: str
    profile_name: str
    content: str
    instructions: Optional[str] = None
    system_prompt: str
    config: Dict[str, Any] = {}


class StudioGenerationOutput(CommandOutput):
    success: bool
    artifact_id: Optional[str] = None
    processing_time: float = 0.0
    error_message: Optional[str] = None


async def _load_artifact(artifact_id: str) -> StudioArtifact:
    artifact = await StudioArtifact.get(artifact_id)
    if not artifact:
        raise ValueError(f"Studio artifact '{artifact_id}' not found")
    return artifact


# --------------------------------------------------------------------------
# ReportWriter — LLM markdown
# --------------------------------------------------------------------------
@command("generate_report", app="open_notebook", retry={"max_attempts": 1})
async def generate_report_command(
    input_data: StudioGenerationInput,
) -> StudioGenerationOutput:
    """ReportWriter agent: generate a markdown report."""
    start = time.time()
    try:
        artifact = await _load_artifact(input_data.artifact_id)
        markdown = await _run_llm(
            input_data.system_prompt, input_data.content, input_data.instructions
        )
        artifact.content = markdown
        await artifact.save()
        logger.info(f"ReportWriter done: {artifact.id}")
        return StudioGenerationOutput(
            success=True,
            artifact_id=str(artifact.id),
            processing_time=time.time() - start,
        )
    except ValueError:
        raise
    except Exception as e:
        logger.error(f"Report generation failed: {e}")
        logger.exception(e)
        raise RuntimeError(str(e)) from e


# --------------------------------------------------------------------------
# QuizMaster — LLM markdown quiz (questions + answers + explanations)
# --------------------------------------------------------------------------
@command("generate_quiz", app="open_notebook", retry={"max_attempts": 1})
async def generate_quiz_command(
    input_data: StudioGenerationInput,
) -> StudioGenerationOutput:
    """QuizMaster agent: generate a markdown quiz with answers."""
    start = time.time()
    try:
        artifact = await _load_artifact(input_data.artifact_id)
        markdown = await _run_llm(
            input_data.system_prompt, input_data.content, input_data.instructions
        )
        artifact.content = markdown
        await artifact.save()
        logger.info(f"QuizMaster done: {artifact.id}")
        return StudioGenerationOutput(
            success=True,
            artifact_id=str(artifact.id),
            processing_time=time.time() - start,
        )
    except ValueError:
        raise
    except Exception as e:
        logger.error(f"Quiz generation failed: {e}")
        logger.exception(e)
        raise RuntimeError(str(e)) from e


# --------------------------------------------------------------------------
# MindMapArchitect — LLM mermaid
# --------------------------------------------------------------------------
@command("generate_mindmap", app="open_notebook", retry={"max_attempts": 1})
async def generate_mindmap_command(
    input_data: StudioGenerationInput,
) -> StudioGenerationOutput:
    """MindMapArchitect agent: generate Mermaid mind-map text."""
    start = time.time()
    try:
        artifact = await _load_artifact(input_data.artifact_id)
        mermaid = await _run_llm(
            input_data.system_prompt, input_data.content, input_data.instructions
        )
        artifact.content = repair_mermaid(_strip_code_fence(mermaid))
        await artifact.save()
        logger.info(f"MindMapArchitect done: {artifact.id}")
        return StudioGenerationOutput(
            success=True,
            artifact_id=str(artifact.id),
            processing_time=time.time() - start,
        )
    except ValueError:
        raise
    except Exception as e:
        logger.error(f"Mindmap generation failed: {e}")
        logger.exception(e)
        raise RuntimeError(str(e)) from e


# --------------------------------------------------------------------------
# SlideComposer — LLM slide prompts -> Doubao images -> stitched .pptx
# --------------------------------------------------------------------------
def _parse_prompt_list(raw: str, limit: int) -> List[str]:
    """Parse the LLM's JSON array of image prompts; tolerate stray fences."""
    text = _strip_code_fence(raw)
    prompts: List[str] = []
    try:
        parsed = json.loads(text)
        if isinstance(parsed, list):
            prompts = [str(p).strip() for p in parsed if str(p).strip()]
    except json.JSONDecodeError:
        # Fallback: treat each non-empty line as a prompt.
        prompts = [ln.strip("-• \t") for ln in text.splitlines() if ln.strip()]
    if not prompts:
        raise ValueError("LLM did not return any image prompts")
    return prompts[:limit]


def _build_pptx(image_paths: List[str], dest: Path) -> None:
    """Stitch slide images into a single 16:9 .pptx, one image per slide.

    Each image is scaled to fit the slide while preserving aspect ratio and
    centered on a blank layout (a deck of full-bleed visual slides).
    """
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    # 16:9 widescreen deck.
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    blank_layout = prs.slide_layouts[6]  # fully blank

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        with Image.open(img_path) as im:
            img_w, img_h = im.size
        # Scale to fit (contain), preserving aspect ratio, then center.
        scale = min(slide_w / img_w, slide_h / img_h)
        draw_w = int(img_w * scale)
        draw_h = int(img_h * scale)
        left = int((slide_w - draw_w) / 2)
        top = int((slide_h - draw_h) / 2)
        slide.shapes.add_picture(
            img_path, left, top, width=draw_w, height=draw_h
        )

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))


@command("generate_ppt", app="open_notebook", retry={"max_attempts": 1})
async def generate_ppt_command(
    input_data: StudioGenerationInput,
) -> StudioGenerationOutput:
    """SlideComposer agent: LLM -> N slide prompts -> Doubao images -> .pptx."""
    from open_notebook.ai.doubao import DoubaoImageClient

    start = time.time()
    try:
        artifact = await _load_artifact(input_data.artifact_id)
        num_images = int(input_data.config.get("num_images", 4))
        size = input_data.config.get("size", "1280x720")

        raw = await _run_llm(
            input_data.system_prompt, input_data.content, input_data.instructions
        )
        prompts = _parse_prompt_list(raw, num_images)

        client = DoubaoImageClient()
        out_dir = _artifact_dir(input_data.artifact_id)
        image_paths: List[str] = []
        for idx, prompt in enumerate(prompts):
            result = client.generate(prompt, size=size, response_format="url")
            if not result.url:
                logger.warning(f"Slide image {idx} returned no URL, skipping")
                continue
            dest = out_dir / f"img_{idx}.png"
            _download(result.url, dest)
            image_paths.append(str(dest))

        if not image_paths:
            raise RuntimeError("No slide images were generated")

        # Stitch the slide images into a single .pptx deck.
        pptx_path = out_dir / "slides.pptx"
        _build_pptx(image_paths, pptx_path)

        # file_paths: slide images first (for gallery preview), deck last.
        artifact.file_paths = image_paths + [str(pptx_path)]
        await artifact.save()
        logger.info(
            f"SlideComposer done: {artifact.id} "
            f"({len(image_paths)} slides -> {pptx_path.name})"
        )
        return StudioGenerationOutput(
            success=True,
            artifact_id=str(artifact.id),
            processing_time=time.time() - start,
        )
    except ValueError:
        raise
    except Exception as e:
        logger.error(f"PPT generation failed: {e}")
        logger.exception(e)
        raise RuntimeError(str(e)) from e


# --------------------------------------------------------------------------
# VideoDirector — LLM prompt -> Doubao Seedance mp4
# --------------------------------------------------------------------------
@command("generate_video", app="open_notebook", retry={"max_attempts": 1})
async def generate_video_command(
    input_data: StudioGenerationInput,
) -> StudioGenerationOutput:
    """VideoDirector agent: LLM video prompt -> Doubao Seedance mp4."""
    from open_notebook.ai.doubao import DoubaoVideoClient

    start = time.time()
    try:
        artifact = await _load_artifact(input_data.artifact_id)
        video_prompt = await _run_llm(
            input_data.system_prompt, input_data.content, input_data.instructions
        )
        video_prompt = _strip_code_fence(video_prompt)

        client = DoubaoVideoClient()
        task_id = client.create_task(
            video_prompt,
            resolution=input_data.config.get("resolution"),
            ratio=input_data.config.get("ratio"),
            duration=input_data.config.get("duration"),
        )
        result = client.wait(task_id, max_wait_seconds=600.0)

        if not result.video_url:
            raise RuntimeError("Doubao video task returned no video URL")

        dest = _artifact_dir(input_data.artifact_id) / "video.mp4"
        _download(result.video_url, dest)

        # Keep the generated script as content for reference / accessibility.
        artifact.content = video_prompt
        artifact.file_paths = [str(dest)]
        await artifact.save()
        logger.info(f"VideoDirector done: {artifact.id} -> {dest}")
        return StudioGenerationOutput(
            success=True,
            artifact_id=str(artifact.id),
            processing_time=time.time() - start,
        )
    except ValueError:
        raise
    except Exception as e:
        logger.error(f"Video generation failed: {e}")
        logger.exception(e)
        raise RuntimeError(str(e)) from e


# Map resource_type -> command name, used by the service layer.
COMMAND_BY_TYPE = {
    "report": "generate_report",
    "quiz": "generate_quiz",
    "mindmap": "generate_mindmap",
    "ppt": "generate_ppt",
    "video": "generate_video",
}
